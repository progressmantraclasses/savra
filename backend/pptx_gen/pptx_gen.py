from __future__ import annotations

import logging
import os
import tempfile
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from models import PresentationContent

logger = logging.getLogger(__name__)

_SLIDE_W = 13.333
_SLIDE_H = 7.5

_DARK_BG = RGBColor(0x0A, 0x1F, 0x0A)
_MID_BG = RGBColor(0x14, 0x53, 0x2D)
_CARD = RGBColor(0x1A, 0x3A, 0x1A)
_CARD_2 = RGBColor(0x1F, 0x4D, 0x2B)
_ACCENT1 = RGBColor(0x4A, 0xDE, 0x80)
_ACCENT2 = RGBColor(0xFD, 0xE6, 0x8A)
_ACCENT3 = RGBColor(0x7D, 0xD3, 0xFC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_MUTED = RGBColor(0x94, 0xA3, 0xB8)


def _create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)
    return prs


def _add_full_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _shape_text_frame(shape, margin: float = 0.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def _add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _shape_text_frame(box, margin=margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _hex_to_rgb(hex_code: str | None, fallback: RGBColor) -> RGBColor:
    if not hex_code:
        return fallback
    clean = hex_code.strip().lstrip("#")
    if len(clean) != 6:
        return fallback
    try:
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))
    except ValueError:
        return fallback


def _resolve_palette(content: PresentationContent) -> list[RGBColor]:
    base = [_ACCENT1, _ACCENT2, _ACCENT3]
    if not content.palette:
        return base

    parsed: list[RGBColor] = []
    for index, color in enumerate(content.palette[:5]):
        parsed.append(_hex_to_rgb(color, base[index % len(base)]))
    return parsed or base


def _slide_accent(slide_data, index: int, palette: list[RGBColor]) -> RGBColor:
    default = palette[index % len(palette)]
    return _hex_to_rgb(slide_data.accent_color, default)


def _add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    border_color: RGBColor | None = None,
    border_pt: float = 1.0,
    radius: bool = True,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_pt)
    return card


def _add_title_slide(prs: Presentation, content: PresentationContent, job_id: str, palette: list[RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_full_bg(slide, _DARK_BG)

    _add_card(slide, 7.55, -0.55, 4.6, 4.6, _MID_BG, radius=True)
    _add_card(slide, -0.45, 5.55, 2.0, 2.0, _MID_BG, radius=True)
    _add_card(slide, 10.65, 5.65, 3.2, 2.1, _CARD_2, radius=True)

    _add_textbox(slide, "SAVRA GENERATOR", 0.9, 0.65, 3.5, 0.25, 12, palette[0], bold=True)
    _add_textbox(slide, content.title, 0.7, 1.65, 11.8, 0.9, 34, palette[0], bold=True, align=PP_ALIGN.CENTER)
    subtitle = f"{len(content.slides)} slides · Job {job_id[:8]}"
    _add_textbox(slide, subtitle, 1.75, 2.72, 9.8, 0.3, 15, _WHITE, align=PP_ALIGN.CENTER)
    if content.color_theme:
        _add_textbox(slide, f"Theme: {content.color_theme}", 4.6, 0.68, 4.2, 0.25, 11, palette[1], align=PP_ALIGN.RIGHT)
    _add_textbox(
        slide,
        "Modern layout, strong contrast, and clean section cards.",
        2.0,
        3.22,
        9.4,
        0.3,
        13,
        _MUTED,
        align=PP_ALIGN.CENTER,
    )

    tag = _add_card(slide, 3.35, 6.25, 6.6, 0.55, _CARD, border_color=palette[0], border_pt=1.0)
    tf = _shape_text_frame(tag, margin=0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Generated automatically from your frontend request"
    run.font.size = Pt(11)
    run.font.color.rgb = _MUTED


def _add_section_header(slide, index: int, total: int, heading: str, accent: RGBColor) -> None:
    _add_full_bg(slide, _DARK_BG)
    _add_card(slide, 9.7, -0.7, 4.2, 2.9, _MID_BG, radius=True)
    _add_card(slide, -0.8, 5.7, 3.6, 2.2, _CARD_2, radius=True)
    breadcrumb_color = _MUTED
    title_color = _WHITE
    _add_textbox(slide, f"SLIDE {index:02d} / {total:02d}", 0.52, 0.18, 2.0, 0.22, 11, breadcrumb_color, bold=True)
    _add_textbox(slide, heading, 0.52, 0.52, 10.8, 0.52, 28, title_color, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(1.08), Inches(2.4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()


def _add_bullets_card(slide, bullets: Iterable[str], accent: RGBColor) -> None:
    card = _add_card(slide, 0.58, 1.55, 7.2, 4.85, _CARD, border_color=accent, border_pt=1.25)
    tf = _shape_text_frame(card, margin=0.2)
    tf.clear()

    heading_para = tf.paragraphs[0]
    heading_para.alignment = PP_ALIGN.LEFT
    run = heading_para.add_run()
    run.text = "Key points"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = accent

    items = [b for b in bullets if b][:6]
    if not items:
        items = ["Core concept", "Practical example", "Important takeaway"]

    for bullet in items:
        para = tf.add_paragraph()
        para.level = 0
        para.space_before = Pt(6)
        para.space_after = Pt(0)
        para.text = f"• {bullet}"
        for br in para.runs:
            br.font.size = Pt(18)
            br.font.color.rgb = _WHITE
            br.font.name = "Calibri"


def _add_image_placeholder(
    slide,
    hint: str | None,
    accent: RGBColor,
    *,
    y: float = 1.55,
    h: float = 4.85,
) -> None:
    panel = _add_card(slide, 8.05, y, 4.7, h, _CARD_2, border_color=accent, border_pt=1.25)
    tf = _shape_text_frame(panel, margin=0.14)
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "IMAGE SUGGESTION"
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = accent

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = hint or "Use a high-quality contextual image"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "(auto visual placeholder)"
    r3.font.size = Pt(10)
    r3.font.color.rgb = _MUTED


def _add_stats_panel(slide, stats, accent: RGBColor) -> None:
    panel = _add_card(slide, 8.05, 1.55, 4.7, 2.35, _CARD_2, border_color=accent, border_pt=1.25)
    tf = _shape_text_frame(panel, margin=0.12)
    tf.clear()

    title = tf.paragraphs[0]
    title.alignment = PP_ALIGN.LEFT
    tr = title.add_run()
    tr.text = "Key stats"
    tr.font.size = Pt(11)
    tr.font.bold = True
    tr.font.color.rgb = accent

    rows = stats[:4] if stats else []
    if not rows:
        rows = [
            {"label": "Coverage", "value": "100%"},
            {"label": "Focus", "value": "High"},
        ]

    for item in rows:
        label = getattr(item, "label", None) or item.get("label", "Metric")
        value = getattr(item, "value", None) or item.get("value", "N/A")
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        rr = p.add_run()
        rr.text = f"{label}: {value}"
        rr.font.size = Pt(14)
        rr.font.bold = True
        rr.font.color.rgb = _WHITE


def _add_formula_panel(slide, formula: str | None, accent: RGBColor) -> None:
    panel = _add_card(slide, 8.05, 4.05, 4.7, 2.35, _CARD_2, border_color=accent, border_pt=1.25)
    tf = _shape_text_frame(panel, margin=0.12)
    tf.clear()

    title = tf.paragraphs[0]
    title.alignment = PP_ALIGN.LEFT
    tr = title.add_run()
    tr.text = "Formula"
    tr.font.size = Pt(11)
    tr.font.bold = True
    tr.font.color.rgb = accent

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    rr = p.add_run()
    rr.text = formula or "No formula required for this concept"
    rr.font.size = Pt(20 if formula else 13)
    rr.font.bold = True if formula else False
    rr.font.color.rgb = _WHITE


def _add_smart_slide(prs: Presentation, slide_data, index: int, total: int, palette: list[RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _slide_accent(slide_data, index, palette)
    _add_section_header(slide, index, total, slide_data.heading, accent)

    _add_bullets_card(slide, slide_data.bullets, accent)

    slide_type = slide_data.slide_type
    if slide_type == "stats":
        _add_stats_panel(slide, slide_data.stats, accent)
        _add_image_placeholder(slide, slide_data.image_hint, accent, y=4.05, h=2.35)
    elif slide_type == "formula":
        _add_formula_panel(slide, slide_data.formula, accent)
        _add_image_placeholder(slide, slide_data.image_hint, accent, y=1.55, h=2.35)
    elif slide_type == "mixed":
        _add_stats_panel(slide, slide_data.stats, accent)
        _add_formula_panel(slide, slide_data.formula, accent)
    elif slide_type == "image_focus":
        _add_image_placeholder(slide, slide_data.image_hint, accent)
    else:
        _add_stats_panel(slide, slide_data.stats, accent)

    # Decorative timeline dots make the slide look less plain.
    dot_y = 1.72
    for _ in slide_data.bullets[:5] or [1, 2, 3]:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.95), Inches(dot_y), Inches(0.07), Inches(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        dot_y += 0.82

    note = slide.shapes.add_textbox(Inches(8.2), Inches(6.0), Inches(4.35), Inches(0.38))
    ntf = _shape_text_frame(note, margin=0.05)
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    nr = np.add_run()
    nr.text = slide_data.speaker_note or f"{slide_type.replace('_', ' ').title()} layout"
    nr.font.size = Pt(11)
    nr.font.color.rgb = _MUTED


def generate_pptx(content: PresentationContent, job_id: str, template_path: str) -> str:
    """Generate a Savra-styled .pptx file from structured content."""
    prs = _create_presentation()

    if os.path.exists(template_path):
        try:
            template = Presentation(template_path)
            if template.slide_width and template.slide_height:
                prs.slide_width = template.slide_width
                prs.slide_height = template.slide_height
        except Exception:
            logger.debug("Template load skipped: %s", template_path, exc_info=True)

    palette = _resolve_palette(content)
    _add_title_slide(prs, content, job_id, palette)

    total = max(1, len(content.slides))
    for index, slide_data in enumerate(content.slides, start=1):
        _add_smart_slide(prs, slide_data, index, total, palette)

    output_path = os.path.join(tempfile.gettempdir(), f"{job_id}.pptx")
    prs.save(output_path)
    logger.info("Savra PPTX saved to %s (%d slides)", output_path, len(content.slides) + 1)
    return output_path
