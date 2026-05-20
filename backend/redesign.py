"""
Savra-style PPT redesign: Photosynthesis Class 7
Run: python redesign.py
Output: photosynthesis_savra.pptx
"""

from __future__ import annotations

from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0A, 0x1F, 0x0A)
MID_CARD = RGBColor(0x1A, 0x3A, 0x1A)
DARK_CARD = RGBColor(0x14, 0x53, 0x2D)
ACCENT1 = RGBColor(0x4A, 0xDE, 0x80)
ACCENT2 = RGBColor(0xFD, 0xE6, 0x8A)
ACCENT3 = RGBColor(0x7D, 0xD3, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG = RGBColor(0xF0, 0xFD, 0xF4)
DARK_TEXT = RGBColor(0x0A, 0x1F, 0x0A)
MID_TEXT = RGBColor(0x47, 0x55, 0x69)
LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)
PALE_GREEN = RGBColor(0xE2, 0xF8, 0xE8)

FONT = "Calibri"

# ── Helpers ──────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _configure_text_frame(tf, margin=0.0) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    font_name: str = FONT,
    margin: float = 0.0,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    _configure_text_frame(tf, margin=margin)
    tf.vertical_anchor = valign
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    border_color: RGBColor | None = None,
    border_pt: float = 1.0,
    shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    return shape


def add_text_to_shape(shape, paragraphs: Sequence[dict], margin: float = 0.15) -> None:
    tf = shape.text_frame
    _configure_text_frame(tf, margin=margin)
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for index, spec in enumerate(paragraphs):
        if index == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(spec.get("space_after", 0))
        p.space_before = Pt(spec.get("space_before", 0))
        run = p.add_run()
        run.text = spec.get("text", "")
        run.font.size = Pt(spec.get("size", 14))
        run.font.color.rgb = spec.get("color", WHITE)
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.name = spec.get("font_name", FONT)


def add_breadcrumb(slide, text: str, color: RGBColor = MUTED) -> None:
    add_textbox(slide, text, 0.5, 0.18, 4.8, 0.22, 11, color, bold=True)


def add_title(slide, text: str, color: RGBColor = WHITE, size: int = 36, x: float = 0.5, y: float = 0.48, w: float = 9.0) -> None:
    add_textbox(slide, text, x, y, w, 0.7, size, color, bold=True)


def add_rule(slide, x: float, y: float, w: float, color: RGBColor, h: float = 0.04) -> None:
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def add_circle(slide, x: float, y: float, d: float, color: RGBColor) -> None:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()


def add_pill(slide, x: float, y: float, w: float, h: float, text: str, fill: RGBColor, color: RGBColor = MUTED) -> None:
    pill = add_card(slide, x, y, w, h, fill)
    add_text_to_shape(
        pill,
        [{"text": text, "size": 13, "color": color, "bold": False, "align": PP_ALIGN.CENTER}],
        margin=0.10,
    )


def add_row_card(slide, x: float, y: float, w: float, h: float, accent: RGBColor, title: str, body: str, dark: bool = False):
    base = add_card(slide, x, y, w, h, LIGHT_CARD if not dark else MID_CARD, border_color=accent, border_pt=1.5)
    add_card(slide, x, y, 0.08, h, accent, shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_text_to_shape(
        base,
        [
            {"text": title, "size": 13, "color": accent, "bold": True},
            {"text": body, "size": 13 if not dark else 14, "color": MID_TEXT if not dark else WHITE, "bold": False},
        ],
        margin=0.16,
    )
    return base


def add_small_stat(slide, x: float, y: float, w: float, h: float, label: str, value: str, accent: RGBColor, dark: bool = True):
    card = add_card(slide, x, y, w, h, MID_CARD if dark else LIGHT_CARD, border_color=accent, border_pt=1.25)
    add_text_to_shape(
        card,
        [
            {"text": label, "size": 11, "color": accent, "bold": True},
            {"text": value, "size": 22, "color": WHITE if dark else DARK_TEXT, "bold": True},
        ],
        margin=0.14,
    )
    return card


def add_bulleted_card(slide, x: float, y: float, w: float, h: float, fill: RGBColor, border_color: RGBColor, title: str, body: str, title_color: RGBColor, body_color: RGBColor, body_size: int = 14):
    card = add_card(slide, x, y, w, h, fill, border_color=border_color, border_pt=1.25)
    add_text_to_shape(
        card,
        [
            {"text": title, "size": 12, "color": title_color, "bold": True},
            {"text": body, "size": body_size, "color": body_color, "bold": False},
        ],
        margin=0.15,
    )
    return card


# ── Slide builders ───────────────────────────────────────────────────────────
def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    add_circle(slide, 7.75, -0.65, 4.1, MID_CARD)
    add_circle(slide, -0.55, 5.55, 2.2, MID_CARD)

    add_textbox(slide, "BIOLOGY • CLASS 7", 0.0, 1.2, 10.0, 0.25, 12, ACCENT1, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "PHOTOSYNTHESIS", 0.55, 1.95, 8.9, 0.75, 52, ACCENT1, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "The process that powers all life on Earth", 1.3, 2.9, 7.4, 0.35, 18, WHITE, align=PP_ALIGN.CENTER)
    add_rule(slide, 3.0, 3.42, 4.0, ACCENT1, h=0.04)

    tag = add_card(slide, 2.05, 6.4, 5.9, 0.52, MID_CARD, border_color=ACCENT1, border_pt=1.0)
    add_text_to_shape(
        tag,
        [{"text": "CBSE • Class 7 • Life Sciences", "size": 11, "color": MUTED, "bold": False, "align": PP_ALIGN.CENTER}],
        margin=0.08,
    )


def slide_2_what_is(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_breadcrumb(slide, "FUNDAMENTALS")
    add_title(slide, "WHAT IS PHOTOSYNTHESIS?", WHITE, 36)

    left = add_card(slide, 0.5, 1.6, 5.1, 3.25, MID_CARD, border_color=ACCENT1, border_pt=1.25)
    add_text_to_shape(
        left,
        [
            {"text": "DEFINITION", "size": 12, "color": ACCENT1, "bold": True},
            {"text": "A process by which plants use sunlight, water, and CO₂ to produce glucose and oxygen.", "size": 17, "color": WHITE, "bold": False},
        ],
        margin=0.18,
    )

    right = add_card(slide, 5.85, 1.6, 3.65, 3.25, MID_CARD, border_color=ACCENT2, border_pt=1.25)
    add_text_to_shape(
        right,
        [
            {"text": "THE EQUATION", "size": 12, "color": ACCENT2, "bold": True, "align": PP_ALIGN.CENTER},
            {"text": "6CO₂ + 6H₂O → C₆H₁₂O₆ + 6O₂", "size": 17, "color": ACCENT2, "bold": True, "align": PP_ALIGN.CENTER},
            {"text": "Sunlight + Chlorophyll", "size": 12, "color": MUTED, "bold": False, "align": PP_ALIGN.CENTER},
        ],
        margin=0.18,
    )

    insight = add_card(slide, 0.5, 5.35, 9.0, 1.0, DARK_CARD, border_color=ACCENT1, border_pt=1.0)
    add_text_to_shape(
        insight,
        [
            {"text": "Key Insight:", "size": 13, "color": ACCENT1, "bold": True},
            {"text": "Plants are the only living things that make their own food.", "size": 14, "color": ACCENT1, "italic": True},
        ],
        margin=0.16,
    )


def slide_3_oxygen(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_breadcrumb(slide, "WHY IT MATTERS", color=MUTED)
    add_title(slide, "IMPORTANCE OF OXYGEN", DARK_TEXT, 36)

    card_w = 2.85
    gap = 0.22
    x0 = 0.5
    y = 1.65
    cards = [
        ("RESPIRATION", "Essential for human and animal respiration — every breath depends on plant photosynthesis."),
        ("ENERGY", "Oxygen is needed to burn glucose for cellular energy in all aerobic organisms."),
        ("LIFE SUPPORT", "Supports all cellular activities — without it, complex life cannot exist."),
    ]
    for idx, (title, body) in enumerate(cards):
        x = x0 + idx * (card_w + gap)
        add_bulleted_card(slide, x, y, card_w, 2.7, LIGHT_CARD, ACCENT1, title, body, ACCENT1, MID_TEXT, body_size=14)

    add_textbox(slide, "O₂", 6.95, 4.98, 2.25, 1.55, 84, RGBColor(0xC6, 0xF6, 0xD5), bold=True, align=PP_ALIGN.RIGHT)


def slide_4_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_breadcrumb(slide, "INPUTS")
    add_title(slide, "WHAT DOES A PLANT NEED?", WHITE, 36)

    entries = [
        (0.5, 1.55, 4.2, 1.45, ACCENT2, "SUNLIGHT", "Light energy from the sun drives the entire reaction"),
        (5.1, 1.55, 4.4, 1.45, ACCENT3, "WATER", "Absorbed through roots, travels up the stem to leaves"),
        (0.5, 3.15, 4.2, 1.45, ACCENT1, "CO₂", "Carbon dioxide enters through stomata on leaf surfaces"),
        (5.1, 3.15, 4.4, 1.45, ACCENT1, "CHLOROPHYLL", "The green pigment in chloroplasts that captures light energy"),
    ]
    for x, y, w, h, accent, label, body in entries:
        card = add_card(slide, x, y, w, h, MID_CARD, border_color=accent, border_pt=1.25)
        add_card(slide, x, y, 0.12, h, accent, shape_type=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        add_text_to_shape(
            card,
            [
                {"text": label, "size": 15, "color": accent, "bold": True},
                {"text": body, "size": 13, "color": WHITE, "bold": False},
            ],
            margin=0.16,
        )


def slide_5_plant_parts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_breadcrumb(slide, "PLANT ANATOMY", color=MUTED)
    add_title(slide, "PARTS INVOLVED IN PHOTOSYNTHESIS", DARK_TEXT, 34)

    rows = [
        ("LEAVES", "Primary site; flat surface maximizes light absorption", ACCENT1),
        ("CHLOROPLASTS", "Organelles containing chlorophyll; actual site of reaction", ACCENT3),
        ("STOMATA", "Tiny pores on leaf underside; CO₂ enters, O₂ exits", ACCENT2),
        ("ROOTS", "Absorb water and minerals from soil for transport upward", ACCENT1),
    ]

    for index, (label, body, accent) in enumerate(rows):
        y = 1.7 + index * 1.0
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.62), Inches(y + 0.2), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        add_textbox(slide, label, 0.9, y, 1.75, 0.25, 16, DARK_TEXT, bold=True)
        add_textbox(slide, body, 2.85, y, 2.9, 0.5, 13, MID_TEXT)

    fact = add_card(slide, 6.25, 1.75, 3.2, 4.15, DARK_CARD, border_color=ACCENT1, border_pt=1.25)
    add_text_to_shape(
        fact,
        [
            {"text": "DID YOU KNOW?", "size": 12, "color": ACCENT1, "bold": True, "align": PP_ALIGN.CENTER},
            {"text": "A single leaf has thousands of stomata per square centimeter.", "size": 21, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
        ],
        margin=0.16,
    )


def slide_6_byproducts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_breadcrumb(slide, "OUTPUTS")
    add_title(slide, "WHAT DOES PHOTOSYNTHESIS PRODUCE?", WHITE, 34)

    add_small_stat(slide, 0.6, 1.75, 4.15, 3.85, "GLUCOSE", "C₆H₁₂O₆", ACCENT2, dark=True)
    add_textbox(slide, "Used for energy, stored as starch in roots and fruits", 0.86, 4.45, 3.65, 0.6, 13, MUTED)

    add_small_stat(slide, 5.25, 1.75, 4.15, 3.85, "OXYGEN", "O₂", ACCENT1, dark=True)
    add_textbox(slide, "Released into atmosphere — the air we breathe", 5.52, 4.45, 3.65, 0.6, 13, MUTED)


def slide_7_water(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_breadcrumb(slide, "WATER'S ROLE", color=MUTED)
    add_title(slide, "WHY WATER IS ESSENTIAL", DARK_TEXT, 36)

    bar_y = 1.7
    bar_specs = [
        ("PHOTOSYNTHESIS", "Water is split during the light reaction to release oxygen and hydrogen"),
        ("TRANSPIRATION", "Evaporation of water regulates leaf temperature and drives upward flow"),
        ("TURGOR PRESSURE", "Water pressure inside cells keeps leaves firm and upright"),
        ("NUTRIENT TRANSPORT", "Minerals dissolved in water travel from roots to all plant parts"),
    ]
    for idx, (label, body) in enumerate(bar_specs):
        y = bar_y + idx * 0.95
        card = add_card(slide, 0.55, y, 8.9, 0.72, LIGHT_CARD, border_color=RGBColor(0xD8, 0xF3, 0xFF), border_pt=0.9)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(y), Inches(0.08), Inches(0.72))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT3
        bar.line.fill.background()
        add_text_to_shape(
            card,
            [
                {"text": label, "size": 13, "color": DARK_TEXT, "bold": True},
                {"text": body, "size": 12, "color": MID_TEXT, "bold": False},
            ],
            margin=0.16,
        )


def slide_8_types(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_breadcrumb(slide, "TYPES")
    add_title(slide, "TYPES OF PHOTOSYNTHESIS", WHITE, 36)

    card_w = 2.95
    gap = 0.2
    specs = [
        ("OXYGENIC", "In green plants with chlorophyll. Produces oxygen as byproduct. Most common type.", ACCENT1),
        ("ANOXYGENIC", "In bacteria with different pigments. Does NOT produce oxygen. Found in sulfur bacteria.", ACCENT2),
        ("CHEMOSYNTHESIS", "In deep-sea microorganisms. Uses chemical energy instead of light. No sunlight needed.", ACCENT3),
    ]

    for idx, (title, body, accent) in enumerate(specs):
        x = 0.5 + idx * (card_w + gap)
        card = add_card(slide, x, 1.7, card_w, 2.95, MID_CARD, border_color=accent, border_pt=1.25)
        add_text_to_shape(
            card,
            [
                {"text": title, "size": 14, "color": accent, "bold": True},
                {"text": body, "size": 13, "color": WHITE, "bold": False},
            ],
            margin=0.16,
        )

    insight = add_card(slide, 0.5, 5.15, 9.0, 1.0, DARK_CARD, border_color=ACCENT1, border_pt=1.0)
    add_text_to_shape(
        insight,
        [
            {"text": "Most life on Earth depends on oxygenic photosynthesis by green plants.", "size": 14, "color": ACCENT1, "italic": True, "align": PP_ALIGN.CENTER},
        ],
        margin=0.12,
    )


def slide_9_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_circle(slide, -0.55, 5.2, 2.5, MID_CARD)
    add_breadcrumb(slide, "KEY TAKEAWAYS")
    add_textbox(slide, "Photosynthesis is the foundation of all life on Earth", 0.65, 1.75, 8.7, 0.75, 36, WHITE, bold=True, align=PP_ALIGN.CENTER)

    pill_y = 3.8
    add_pill(slide, 0.75, pill_y, 2.75, 0.68, "Produces O₂ we breathe", MID_CARD)
    add_pill(slide, 3.6, pill_y, 2.75, 0.68, "Creates glucose for energy", MID_CARD)
    add_pill(slide, 6.45, pill_y, 2.75, 0.68, "Drives all food chains", MID_CARD)

    add_textbox(slide, "Class 7 • CBSE Science • Chapter: Nutrition in Plants", 1.2, 6.25, 7.6, 0.25, 11, MUTED, align=PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────────────────────
def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_1_title(prs)
    slide_2_what_is(prs)
    slide_3_oxygen(prs)
    slide_4_requirements(prs)
    slide_5_plant_parts(prs)
    slide_6_byproducts(prs)
    slide_7_water(prs)
    slide_8_types(prs)
    slide_9_conclusion(prs)

    prs.save("photosynthesis_savra.pptx")
    print("Saved: photosynthesis_savra.pptx")


if __name__ == "__main__":
    main()
